"""PowerPoint slide builder implementation"""

import os
from pathlib import Path
from typing import List, Optional, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
import pandas as pd
from loguru import logger

from ..parser.models import SlideContent, MarkdownElement, ElementType


class SlideBuilder:
    """Build PowerPoint presentations from parsed Markdown content"""
    
    def __init__(self, template_path: Optional[str] = None):
        """Initialize the slide builder
        
        Args:
            template_path: Path to PowerPoint template file
        """
        if template_path and os.path.exists(template_path):
            self.prs = Presentation(template_path)
            logger.info(f"Loaded template from {template_path}")
        else:
            self.prs = Presentation()
            logger.info("Created new presentation without template")
            
        self.rendered_images = {}  # Cache for rendered images
        
    def build(self, slides: List[SlideContent], output_path: str) -> None:
        """Build PowerPoint presentation from slide content
        
        Args:
            slides: List of SlideContent objects
            output_path: Path to save the presentation
        """
        logger.info(f"Building presentation with {len(slides)} slides")
        
        for i, slide_content in enumerate(slides):
            if i == 0 and slide_content.slide_index == 0:
                # Title slide
                self._create_title_slide(slide_content)
            else:
                # Content slide
                self._create_content_slide(slide_content)
                
        # Save presentation
        self.prs.save(output_path)
        logger.info(f"Saved presentation to {output_path}")
        
    def _create_title_slide(self, slide_content: SlideContent):
        """Create a title slide"""
        # Use title slide layout (usually index 0)
        slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        if title_shape:
            title_shape.text = slide_content.title
            
        # Add subtitle if available (from first paragraph)
        if slide_content.elements:
            for element in slide_content.elements:
                if element.type == ElementType.PARAGRAPH:
                    # Look for subtitle placeholder
                    for shape in slide.placeholders:
                        if shape.placeholder_format.idx == 1:  # Subtitle placeholder
                            shape.text = element.content
                            break
                    break
                    
        # Add notes
        if slide_content.notes:
            slide.notes_slide.notes_text_frame.text = slide_content.notes
            
    def _create_content_slide(self, slide_content: SlideContent):
        """Create a content slide"""
        # Use content slide layout (usually index 1)
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        if title_shape:
            title_shape.text = slide_content.title
            
        # Find or create content placeholder
        content_shape = None
        lead_shape = None
        
        # Look for Lead placeholder and content placeholder
        for shape in slide.shapes:
            if hasattr(shape, 'name'):
                if shape.name == "Lead" and slide_content.lead_text:
                    lead_shape = shape
                elif shape.placeholder_format and shape.placeholder_format.idx == 1:
                    content_shape = shape
                    
        # Set lead text if available
        if lead_shape and slide_content.lead_text:
            lead_shape.text = slide_content.lead_text
            
        # Process content elements
        if content_shape:
            text_frame = content_shape.text_frame
            text_frame.clear()  # Clear existing content
            
            for element in slide_content.elements:
                self._add_element_to_slide(slide, element, text_frame)
        else:
            # No content placeholder, add elements directly to slide
            top = Inches(2)
            for element in slide_content.elements:
                top = self._add_element_to_slide_direct(slide, element, top)
                
        # Add notes
        if slide_content.notes:
            slide.notes_slide.notes_text_frame.text = slide_content.notes
            
    def _add_element_to_slide(self, slide, element: MarkdownElement, text_frame):
        """Add an element to existing text frame"""
        if element.type == ElementType.PARAGRAPH:
            p = text_frame.add_paragraph()
            p.text = element.content
            p.level = 0
            
        elif element.type in (ElementType.LIST_UNORDERED, ElementType.LIST_ORDERED):
            for i, item in enumerate(element.children):
                p = text_frame.add_paragraph()
                if element.type == ElementType.LIST_ORDERED:
                    p.text = f"{i+1}. {item.content}"
                else:
                    p.text = f"• {item.content}"
                p.level = 1
                
        elif element.type == ElementType.CODE_BLOCK:
            p = text_frame.add_paragraph()
            p.text = element.content
            p.font.name = "Consolas"
            p.font.size = Pt(10)
            
    def _add_element_to_slide_direct(self, slide, element: MarkdownElement, top: float) -> float:
        """Add an element directly to slide (when no placeholder exists)"""
        left = Inches(0.5)
        width = Inches(9)
        
        if element.type == ElementType.IMAGE:
            # Add image
            img_path = element.content
            if os.path.exists(img_path):
                pic = slide.shapes.add_picture(img_path, left, top, width=width)
                # Maintain aspect ratio
                aspect_ratio = pic.height / pic.width
                pic.width = width
                pic.height = int(width * aspect_ratio)
                return top + pic.height + Inches(0.2)
                
        elif element.type == ElementType.MERMAID:
            # Add Mermaid diagram (assuming it's been rendered)
            if element.content in self.rendered_images:
                img_path = self.rendered_images[element.content]
                pic = slide.shapes.add_picture(img_path, left, top, width=width)
                aspect_ratio = pic.height / pic.width
                pic.width = width
                pic.height = int(width * aspect_ratio)
                return top + pic.height + Inches(0.2)
                
        elif element.type == ElementType.TABLE:
            # Add table
            headers = element.attributes.get("headers", [])
            rows = element.attributes.get("rows", [])
            
            if headers and rows:
                table = slide.shapes.add_table(
                    len(rows) + 1, len(headers), left, top, width, Inches(0.5)
                ).table
                
                # Set headers
                for i, header in enumerate(headers):
                    table.cell(0, i).text = header
                    
                # Set data
                for i, row in enumerate(rows):
                    for j, cell in enumerate(row):
                        table.cell(i + 1, j).text = str(cell)
                        
                return top + Inches(0.5) * (len(rows) + 2)
                
        else:
            # Add text box for other elements
            height = Inches(0.5)
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.word_wrap = True
            
            if element.type == ElementType.PARAGRAPH:
                text_frame.text = element.content
            elif element.type == ElementType.CODE_BLOCK:
                text_frame.text = element.content
                for paragraph in text_frame.paragraphs:
                    paragraph.font.name = "Consolas"
                    paragraph.font.size = Pt(10)
            elif element.type in (ElementType.LIST_UNORDERED, ElementType.LIST_ORDERED):
                text_frame.clear()
                for i, item in enumerate(element.children):
                    p = text_frame.add_paragraph()
                    if element.type == ElementType.LIST_ORDERED:
                        p.text = f"{i+1}. {item.content}"
                    else:
                        p.text = f"• {item.content}"
                        
            # Auto-size the text box
            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            
            return top + height + Inches(0.2)
            
    def add_rendered_image(self, mermaid_code: str, image_path: str):
        """Register a rendered Mermaid image
        
        Args:
            mermaid_code: Original Mermaid code
            image_path: Path to rendered image
        """
        self.rendered_images[mermaid_code] = image_path